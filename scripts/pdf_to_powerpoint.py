import sys
import tempfile
from pathlib import Path

import pymupdf
from pptx import Presentation
from pptx.util import Inches


def convert_pdf_to_powerpoint(
    input_pdf: str,
    output_pptx: str,
    dpi: int = 180,
) -> None:
    input_path = Path(input_pdf)
    output_path = Path(output_pptx)

    if not input_path.exists():
        raise FileNotFoundError(
            f"PDF file not found: {input_path}"
        )

    if input_path.suffix.lower() != ".pdf":
        raise ValueError(
            "The input file must be a PDF."
        )

    output_path.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    pdf_document = pymupdf.open(
        str(input_path)
    )

    if pdf_document.page_count == 0:
        pdf_document.close()

        raise ValueError(
            "The PDF does not contain any pages."
        )

    presentation = Presentation()

    # Remove the default presentation slide size.
    first_page = pdf_document[0]
    first_rect = first_page.rect

    # PowerPoint size based on first PDF page ratio.
    slide_width_inches = 10
    slide_height_inches = (
        slide_width_inches
        * first_rect.height
        / first_rect.width
    )

    presentation.slide_width = Inches(
        slide_width_inches
    )

    presentation.slide_height = Inches(
        slide_height_inches
    )

    blank_slide_layout = (
        presentation.slide_layouts[6]
    )

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        for page_index in range(
            pdf_document.page_count
        ):
            page = pdf_document[page_index]

            pixmap = page.get_pixmap(
                dpi=dpi,
                alpha=False,
            )

            image_path = (
                temp_path
                / f"page_{page_index + 1}.png"
            )

            pixmap.save(
                str(image_path)
            )

            slide = presentation.slides.add_slide(
                blank_slide_layout
            )

            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

    pdf_document.close()

    presentation.save(
        str(output_path)
    )


def main() -> None:
    if len(sys.argv) != 3:
        print(
            "Usage: python pdf_to_powerpoint.py "
            "input.pdf output.pptx",
            file=sys.stderr,
        )

        sys.exit(1)

    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]

    try:
        convert_pdf_to_powerpoint(
            input_pdf,
            output_pptx,
        )

        print(
            "PDF converted to PowerPoint successfully."
        )

    except Exception as error:
        print(
            f"PDF to PowerPoint conversion failed: "
            f"{error}",
            file=sys.stderr,
        )

        sys.exit(1)


if __name__ == "__main__":
    main()